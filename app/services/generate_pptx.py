"""
Generate a PowerPoint from Microsoft Dataverse data using a fixed template.

Dependencies:
    pip install msal requests python-pptx

Configuration:
    Set the following environment variables:
        DATAVERSE_CLIENT_ID   - Azure AD application (client) ID
        DATAVERSE_CLIENT_SECRET - Azure AD application client secret
        DATAVERSE_TENANT_ID   - Azure AD tenant ID
        DATAVERSE_URL         - Dataverse endpoint, e.g. https://<org>.api.crm.dynamics.com/api/data/v9.1/

Usage:
    Call the functions in this module to generate PowerPoint presentations.
"""
import os
import json
from pptx import Presentation
from pptx.util import Pt
from datetime import datetime
import requests
import msal

CLIENT_ID = os.getenv('DATAVERSE_CLIENT_ID')
CLIENT_SECRET = os.getenv('DATAVERSE_CLIENT_SECRET')
TENANT_ID = os.getenv('DATAVERSE_TENANT_ID')
DATAVERSE_URL = os.getenv('DATAVERSE_URL')
DATAVERSE_API_URL = os.getenv('DATAVERSE_API_URL')

PPTX_TEMPLATE = os.getenv('PPTX_TEMPLATE', 'template.pptx')
OUTPUT_PATH = os.path.join(os.getcwd(), "output")

SCOPE = [f"{DATAVERSE_URL.strip('/')}/.default"]

def get_access_token():
    authority = f"https://login.microsoftonline.com/{TENANT_ID}"
    app = msal.ConfidentialClientApplication(
        CLIENT_ID,
        authority=authority,
        client_credential=CLIENT_SECRET
    )
    result = app.acquire_token_for_client(scopes=SCOPE)
    if 'access_token' in result:
        return result['access_token']
    else:
        raise Exception(f"Token acquisition failed: {result.get('error_description')}")

def fetch_data(entity: str, select=None, filter_expr=None):
    token = get_access_token()
    if not token:
        raise Exception("Failed to acquire access token")
    headers = {
        'Authorization': f'Bearer {token}',
        'Accept': 'application/json'
    }
    params = {}
    if select:
        params['$select'] = ','.join(select)
    if filter_expr:
        params['$filter'] = filter_expr

    url = f"{DATAVERSE_API_URL.rstrip('/')}/{entity}"
    resp = requests.get(url, headers=headers, params=params)
    resp.raise_for_status()
    return resp.json().get('value', [])

def iter_cells(table):
    """Helper function to iterate over all cells in a table."""
    for row in table.rows:
        for cell in row.cells:
            yield cell

def set_table_font_size(table, font_size):
    """Set the font size for all cells in a table."""
    for cell in iter_cells(table):
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(font_size)

def create_table(slide, rows, cols, left, top, width, height, font_size, font_bold):
    new_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    new_table = new_shape.table

    for cell in iter_cells(new_table):
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(font_size)
                run.font.bold = font_bold

    return new_table

def process_table_placeholder(table, content, slide):
    placeholder_name = table.cell(0, 0).text[8:-2].strip()
    value = content.get(placeholder_name, [])
    if not isinstance(value, list) or not value:
        table.cell(0, 0).text = "n/a"
        set_table_font_size(table, 11)
        return

    headers = list(value[0].keys())
    rows, cols = len(value) + 1, len(headers)

    left = table._graphic_frame.left
    top = table._graphic_frame.top
    width = table._graphic_frame.width
    height = table._graphic_frame.height
    font_size = table.cell(0, 0).text_frame.paragraphs[0].font.size
    font_bold = table.cell(0, 0).text_frame.paragraphs[0].font.bold

    sp = table._graphic_frame._element
    sp.getparent().remove(sp)

    new_table = create_table(slide, rows, cols, left, top, width, height, font_size, font_bold)

    for col_idx, header in enumerate(headers):
        cell = new_table.cell(0, col_idx)
        cell.text = header

    for row_idx, item in enumerate(value, start=1):
        for col_idx, header in enumerate(headers):
            cell = new_table.cell(row_idx, col_idx)
            cell.text = str(item.get(header, ""))

    set_table_font_size(new_table, 11)

def process_text_placeholders(slide, content):
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                fulltext = ''.join(run.text for run in paragraph.runs)

                for placeholder_name, value in content.items():
                    placeholder_tag = f"{{{{{placeholder_name}}}}}"
                    if placeholder_tag in fulltext:
                        fulltext = fulltext.replace(placeholder_tag, str(value))

                while "{{" in fulltext and "}}" in fulltext:
                    start_idx = fulltext.find("{{")
                    end_idx = fulltext.find("}}", start_idx) + 2
                    placeholder_tag = fulltext[start_idx:end_idx]
                    fulltext = fulltext.replace(placeholder_tag, "n/a")

                for run in paragraph.runs:
                    run.text = ''
                if paragraph.runs:
                    paragraph.runs[0].text = fulltext

def generate_ppt(jobid: str, records: list, template_path='templates/template.pptx', output_filename='output.pptx'):
    print(os.getcwd())
    
    prs = Presentation("app/templates/template.pptx")
    jobdate = datetime.now().strftime('%Y-%m-%d %H:%M:%S')

    for idx, record in enumerate(records):
        content = {}
        if 'jeschro_content' in record:
            try:
                content = json.loads(record['jeschro_content'])
            except json.JSONDecodeError as e:
                continue

        content['jobid'] = jobid
        content['jobdate'] = jobdate

        for slide in prs.slides:
            process_text_placeholders(slide, content)

            for shape in slide.shapes:
                if shape.has_table:
                    process_table_placeholder(shape.table, content, slide)

    if not os.path.exists(OUTPUT_PATH):
        os.makedirs(OUTPUT_PATH)

    prs.save(os.path.join(OUTPUT_PATH, output_filename))
    print(f"Presentation saved as {output_filename}")
    # Post the saved presentation back to Dataverse
    file_path = os.path.join(OUTPUT_PATH, output_filename)
    token = get_access_token()

    # Step 1: Create the record with the relationship to the job
    headers = {
        'Authorization': f'Bearer {token}',
        'Content-Type': 'application/json',
        'Prefer': 'return=representation'
    }
    data = {
        'jeschro_name': output_filename,
        'jeschro_Job@odata.bind': f"/jeschro_jobs({jobid})"
    }
    url = f"{DATAVERSE_API_URL.rstrip('/')}/jeschro_files"
    create_response = requests.post(url, headers=headers, data=json.dumps(data))
    create_response.raise_for_status()

    # Extract the ID of the created record
    created_record_id = create_response.json().get('jeschro_fileid')
    if not created_record_id:
        raise Exception("Failed to retrieve the ID of the created record.")

    # Step 2: Upload the file to the created record
    upload_url = f"{url}({created_record_id})/jeschro_file"
    print(f"Uploading file {output_filename} to {upload_url}")
    with open(file_path, 'rb') as file_data:
        upload_headers = {
            'Authorization': f'Bearer {token}',
            'Content-Type': 'application/octet-stream',
            'x-ms-file-name': output_filename
        }
        upload_response = requests.patch(upload_url, headers=upload_headers, data=file_data)
        upload_response.raise_for_status()

    os.remove(os.path.join(OUTPUT_PATH, output_filename))
    print(f"File {output_filename} successfully uploaded to Dataverse record {created_record_id}.")